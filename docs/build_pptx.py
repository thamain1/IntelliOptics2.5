import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT = r"C:/Dev/intellioptics_2.5/docs/IntelliOptics2.5_Updated.pptx"

NAVY    = (11, 20, 38)
BLUE    = (0, 120, 255)
CYAN    = (0, 201, 255)
WHITE   = (255, 255, 255)
OFFWHT  = (241, 245, 249)
DKGRAY  = (30, 41, 59)
MDGRAY  = (100, 116, 139)
GREEN   = (16, 185, 129)
ORANGE  = (245, 158, 11)
RED     = (239, 68, 68)
DKBLUE  = (0, 40, 80)
MDBLUE  = (0, 60, 120)
LTBLUE  = (0, 80, 160)

def C(t): return RGBColor(*t)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

blank = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(blank)

def bg(sl, color):
    f = sl.background.fill
    f.solid()
    f.fore_color.rgb = C(color)

def rect(sl, l, t, w, h, color):
    s = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = C(color)
    s.line.fill.background()
    return s

def tx(sl, text, l, t, w, h, size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    b = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = b.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = C(color)
    return b

def txb(sl, lines, l, t, w, h, default_size=12, default_color=WHITE, default_bold=False, default_align=PP_ALIGN.LEFT):
    b = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = b.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if isinstance(line, str):
            text, size, bold, color, align = line, default_size, default_bold, default_color, default_align
        else:
            text = line[0]
            size = line[1] if len(line) > 1 else default_size
            bold = line[2] if len(line) > 2 else default_bold
            color = line[3] if len(line) > 3 else default_color
            align = line[4] if len(line) > 4 else default_align
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = C(color)
    return b

def dark_frame(sl):
    bg(sl, NAVY)
    rect(sl, 0, 0, 0.08, 7.5, BLUE)
    rect(sl, 0, 7.1, 13.333, 0.4, CYAN)

def content_frame(sl, title, subtitle=None):
    bg(sl, WHITE)
    rect(sl, 0, 0, 13.333, 1.25, NAVY)
    rect(sl, 0, 0, 0.08, 7.5, BLUE)
    rect(sl, 0, 7.1, 13.333, 0.4, NAVY)
    tx(sl, title, 0.25, 0.08, 11, 0.9, size=26, bold=True, color=WHITE)
    if subtitle:
        tx(sl, subtitle, 0.25, 0.88, 11, 0.4, size=12, color=CYAN)

def card(sl, l, t, w, h, bg_color, title, title_color, bullets, bullet_color=WHITE, title_size=13, bullet_size=11):
    rect(sl, l, t, w, h, bg_color)
    lines = [( title, title_size, True, title_color, PP_ALIGN.LEFT )]
    for b in bullets:
        lines.append((b, bullet_size, False, bullet_color, PP_ALIGN.LEFT))
    txb(sl, lines, l+0.15, t+0.15, w-0.3, h-0.3)

# SLIDE 1
s = slide()
dark_frame(s)
rect(s, 0, 2.8, 6, 0.06, CYAN)
tx(s, "IntelliOptics", 0.35, 1.5, 12, 1.5, size=54, bold=True, color=WHITE)
tx(s, "2.5", 7.3, 1.5, 5, 1.5, size=54, bold=True, color=CYAN)
tx(s, "AI-Powered Visual Intelligence Platform", 0.35, 3.1, 12, 0.9, size=22, color=CYAN)
tx(s, "Haile Hantal  |  Co-Founder, CXO  |  www.4wardmotions.com", 0.35, 4.1, 12, 0.5, size=13, color=(150,170,200))
rect(s, 11.2, 0.9, 1.9, 0.55, BLUE)
tx(s, "VERSION 2.5", 11.25, 0.95, 1.8, 0.45, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# SLIDE 2
s = slide()
content_frame(s, "The Problem We Solve", "Security teams are overwhelmed - not undertooled")
card(s, 0.2, 1.4, 6.3, 5.4, NAVY, "SURVEILLANCE CHALLENGES", CYAN,
    ["50% of cameras experience issues annually - undetected",
     "Misaligned cameras trigger zero alerts when pointing wrong",
     "Manual monitoring fails to scale across hundreds of feeds",
     "Security teams lack time to review hours of footage daily",
     "Critical moments missed between manual check cycles",
     "No early warning system for gradual feed degradation"])
card(s, 6.8, 1.4, 6.3, 5.4, OFFWHT, "THE COST OF DOING NOTHING", BLUE,
    ["Blind spots during critical incidents - no footage to review",
     "Compliance failures from undocumented camera outages",
     "$57,000+ in avoidable labor and equipment costs",
     "Liability exposure when footage is unavailable in court",
     "Reactive security posture - responding, never preventing",
     "IT burdened with manual camera health checks daily"],
    bullet_color=DKGRAY)

# SLIDE 3
s = slide()
dark_frame(s)
rect(s, 0.4, 1.3, 12.5, 2.8, DKBLUE)
txb(s, [
    ('"We are not more footage -', 28, True, WHITE, PP_ALIGN.CENTER),
    ('we are the assurance that the footage exists."', 28, True, WHITE, PP_ALIGN.CENTER),
], 0.5, 1.5, 12.3, 2.3)
tx(s, "IntelliOptics does not just record - it comprehends, validates, and alerts.", 0.5, 4.2, 12.3, 0.6, size=16, italic=True, color=CYAN, align=PP_ALIGN.CENTER)
for x, num, lbl in [(0.4,"376","Detection Classes"),(4.6,"<30s","Edge Response Time"),(8.8,"10fps","Camera Inspection Rate")]:
    rect(s, x, 5.0, 4.0, 1.8, MDBLUE)
    tx(s, num, x+0.1, 5.1, 3.8, 0.9, size=38, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
    tx(s, lbl, x+0.1, 5.95, 3.8, 0.6, size=13, color=WHITE, align=PP_ALIGN.CENTER)

# SLIDE 4
s = slide()
content_frame(s, "What is New in IntelliOptics 2.5", "A major leap in AI detection, language understanding, and operational modules")
card(s, 0.2, 1.4, 6.3, 2.7, BLUE, "YOLOE OPEN-VOCABULARY DETECTION", WHITE,
    ["376-class ONNX inference engine - no Ultralytics dependency",
     "Apache 2.0 licensed - production-safe commercial deployment",
     "Detects objects without retraining on new classes",
     "2x2 tiled detection for dense scenes and overhead views"])
card(s, 6.8, 1.4, 6.3, 2.7, NAVY, "VISION LANGUAGE MODEL (VLM)", CYAN,
    ["Moondream 2B - ask anything about your camera feeds",
     "Natural language queries: Is the gate open?",
     "VLM auto-escalates when YOLOE confidence is low (<0.25)",
     "Endpoints: /vlm/query, /vlm/detect, /vlm/ocr"],
    bullet_color=OFFWHT)
card(s, 0.2, 4.3, 6.3, 2.7, OFFWHT, "SPECIALIZED INTELLIGENCE MODULES", BLUE,
    ["Vehicle ID: license plate OCR + color + vehicle type",
     "Forensic BOLO: search hours of video in minutes",
     "IntelliPark: automated lot monitoring and management",
     "PPE and Safety: personal protective equipment detection"],
    bullet_color=DKGRAY)
card(s, 6.8, 4.3, 6.3, 2.7, DKGRAY, "ENTERPRISE DEPLOYMENT", CYAN,
    ["One-command install: Install-IntelliOptics.ps1",
     "Unified Docker stack - cloud + edge on single network",
     "Pre-filled credentials - live in under a day",
     "17 frontend pages, 19 backend routers, 170+ API endpoints"],
    bullet_color=OFFWHT)

# SLIDE 5
s = slide()
content_frame(s, "IntelliOptics 2.5 Key Features", "Built for enterprise security operations at any scale")
cards_r1 = [
    (NAVY,    "OPEN-VOCAB AI DETECTION",   CYAN,  ["YOLOE 376-class ONNX","Novel prompt VLM fallback","PPE and custom class support","Overhead tiling for dense scenes"]),
    (DKBLUE,  "VISION LANGUAGE MODEL",     CYAN,  ["Ask cameras in plain English","Query / Detect / OCR modes","Moondream 2B - 2025 revision","Auto-escalation pipeline"]),
    (MDBLUE,  "CAMERA HEALTH MONITORING",  CYAN,  ["Twice-daily day/night inspection","Detects: blur, brightness, misalign","10fps inspection capability","Historical health analytics"]),
]
cards_r2 = [
    ((0,50,100), "BEST-IN-CLASS SECURITY",  CYAN,  ["SOC 2 Compliant architecture","AES-256 data encryption","Role-based access control","Full audit trail logging"]),
    ((0,30,70),  "EDGE-FIRST ARCHITECTURE", CYAN,  ["On-site processing - low latency","No cloud dependency for alerts","Custom AI rules per location","Works with existing cameras"]),
    (DKGRAY,     "RAPID DEPLOYMENT",        CYAN,  ["Live in under 1 day","No camera replacement needed","Unified Docker stack","Admin accounts auto-provisioned"]),
]
for i,(bgc,ttl,tc,buls) in enumerate(cards_r1):
    card(s, 0.2+i*4.37, 1.4, 4.2, 2.7, bgc, ttl, tc, buls)
for i,(bgc,ttl,tc,buls) in enumerate(cards_r2):
    card(s, 0.2+i*4.37, 4.3, 4.2, 2.7, bgc, ttl, tc, buls)

# SLIDE 6
s = slide()
content_frame(s, "How IntelliOptics Works", "Edge inference + Cloud intelligence + Human escalation")
steps = [
    (BLUE,       '1\nCAMERA\nFEED'),
    (DKBLUE,     '2\nEDGE\nDETECTION\n(YOLOE)'),
    ((0,55,110), '3\nVLM\nANALYSIS\n(Moondream)'),
    ((0,70,140), '4\nCLOUD\nALERTS &\nLOGGING'),
    ((0,90,175), '5\nDASHBOARD\n& REPORTS'),
]
xs = [0.2, 2.65, 5.1, 7.55, 10.0]
for i,(bgc,lbl) in enumerate(steps):
    rect(s, xs[i], 1.4, 2.2, 2.4, bgc)
    tx(s, lbl, xs[i]+0.1, 1.5, 2.0, 2.2, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 4:
        tx(s, ">>", xs[i]+2.22, 2.1, 0.4, 1.0, size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
notes = [
    ("EDGE LAYER", "YOLOE ONNX runs on-site. Detections in <30s. VLM escalates for low-confidence or novel prompts."),
    ("AI LAYER", "Moondream 2B VLM answers natural language queries. Handles PPE, BOLO, vehicle ID, and OCR."),
    ("CLOUD LAYER", "Unified backend: 170+ endpoints. Real-time alerts, reporting, forensic search, parking management."),
]
nxs = [0.2, 4.6, 9.0]
for i,(ttl,body) in enumerate(notes):
    rect(s, nxs[i], 4.1, 4.1, 2.7, OFFWHT)
    txb(s, [(ttl,13,True,BLUE,PP_ALIGN.LEFT),(body,11,False,DKGRAY,PP_ALIGN.LEFT)], nxs[i]+0.15, 4.25, 3.8, 2.4)

# SLIDE 7
s = slide()
content_frame(s, "Open-Vocabulary AI Detection", "YOLOE - Detect anything you can describe, without retraining")
rect(s, 0.2, 1.4, 6.5, 5.5, NAVY)
txb(s, [
    ("YOLOE CAPABILITIES", 14, True, CYAN, PP_ALIGN.LEFT),
    ("376-class detection vocabulary baked at build time", 11, False, WHITE, PP_ALIGN.LEFT),
    ("Pure ONNX Runtime - no Ultralytics at runtime", 11, False, WHITE, PP_ALIGN.LEFT),
    ("Apache 2.0 licensed - safe for all commercial use", 11, False, WHITE, PP_ALIGN.LEFT),
    ("2x2 tiled detection for parking lots and overhead views", 11, False, WHITE, PP_ALIGN.LEFT),
    ("Cross-tile NMS merging for seamless dense-scene coverage", 11, False, WHITE, PP_ALIGN.LEFT),
    ("Per-prompt VLM fallback: auto-escalates missed detections", 11, False, WHITE, PP_ALIGN.LEFT),
    ("Low-confidence escalation threshold: 0.25 confidence score", 11, False, WHITE, PP_ALIGN.LEFT),
    ("Live 30fps bounding box overlay on frontend canvas", 11, False, WHITE, PP_ALIGN.LEFT),
], 0.35, 1.55, 6.2, 5.2)
card(s, 7.0, 1.4, 6.1, 1.7, BLUE, "STANDARD DETECTION", WHITE,
    ["Detects 376 classes: vehicles, people, objects, signage","Zero configuration - works out of the box"])
card(s, 7.0, 3.3, 6.1, 1.7, DKBLUE, "TILED DETECTION", CYAN,
    ["2x2 grid subdivides frame for high-density scenes","Dramatically improves overhead and parking lot accuracy"],
    bullet_color=OFFWHT)
card(s, 7.0, 5.2, 6.1, 1.7, MDBLUE, "VLM FALLBACK", CYAN,
    ["Fires when YOLOE returns 0 detections OR confidence <0.25","Per-prompt escalation - seamless hybrid AI pipeline"],
    bullet_color=OFFWHT)

# SLIDE 8
s = slide()
content_frame(s, "Vision Language Model - Ask Anything", "Moondream 2B - Natural language intelligence for your camera feeds")
rect(s, 0.2, 1.38, 12.9, 0.75, BLUE)
tx(s, '"Is the gate open?"  |  "Read the license plate."  |  "Is anyone in the restricted zone?"  |  "Is PPE being worn?"',
   0.3, 1.45, 12.7, 0.6, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
card(s, 0.2, 2.3, 6.3, 2.4, NAVY, "/vlm/query", CYAN,
    ["Open-ended conversational queries about any frame",
     "Ask: Describe what you see, Is the area clear?",
     "Returns natural language answer - no training needed",
     "Moondream 2B fp32 - full accuracy, not quantized"],
    bullet_color=OFFWHT)
card(s, 6.8, 2.3, 6.3, 2.4, DKBLUE, "/vlm/detect", CYAN,
    ["Two-pass: presence check + bounding box estimation",
     "Returns bbox coordinates [x1, y1, x2, y2]",
     "Spatial reasoning via natural language prompting",
     "Fallback to full-frame [0,0,1,1] when uncertain"],
    bullet_color=OFFWHT)
card(s, 0.2, 4.85, 6.3, 2.4, (0,20,40), "/vlm/ocr", CYAN,
    ["Read all visible text from any camera frame",
     "License plates, signage, documents, display screens",
     "Powers the Vehicle ID plate recognition pipeline",
     "No separate OCR model required"],
    bullet_color=OFFWHT)
card(s, 6.8, 4.85, 6.3, 2.4, (11,50,90), "VLM AUTO-ESCALATION", CYAN,
    ["Fires when YOLOE returns 0 detections for a prompt",
     "Fires when max detection confidence < 0.25",
     "Per-prompt: each class can escalate independently",
     "Transparent to end user - seamless hybrid pipeline"],
    bullet_color=OFFWHT)

# SLIDE 9
s = slide()
content_frame(s, "Custom Edge Intelligence for Your Business", "Real-time decisions made on-site - no cloud round-trip required")
card(s, 0.2, 1.4, 6.3, 2.7, NAVY, "ON-SITE PROCESSING", CYAN,
    ["YOLOE inference runs entirely at the edge","Sensitive footage never leaves your premises",
     "Real-time decisions - latency under 30 seconds","Operates independently during internet outages"],
    bullet_color=OFFWHT)
card(s, 6.8, 1.4, 6.3, 2.7, DKBLUE, "CUSTOM AI RULES", CYAN,
    ["Define exactly what to detect per camera zone","Configure alert thresholds and escalation logic",
     "Build visual rules aligned to your environment","No ML expertise required - natural language config"],
    bullet_color=OFFWHT)
card(s, 0.2, 4.3, 6.3, 2.7, MDBLUE, "TARGETED ALERTS & METRICS", CYAN,
    ["Get notified only on what matters to your ops","Per-camera, per-zone, per-class alerting rules",
     "SMS, email, and webhook delivery (SendGrid/Twilio)","Define your own KPIs for informed decisions"],
    bullet_color=OFFWHT)
card(s, 6.8, 4.3, 6.3, 2.7, DKGRAY, "FLEXIBLE FIT", CYAN,
    ["Works with any existing IP camera system","No hardware replacement - add AI as a layer",
     "Scales from 1 camera to enterprise-wide deployment","Edge + cloud on a single unified Docker stack"],
    bullet_color=OFFWHT)

# SLIDE 10
s = slide()
content_frame(s, "Specialized Intelligence Modules", "Plug-in capabilities for industry-specific operations")
modules = [
    (DKBLUE, "VEHICLE ID",
     ["License plate OCR","Vehicle color detection","Vehicle type classification","(car, truck, SUV, motorcycle)","","Use cases:","  Parking access control","  Fleet monitoring","  Gate & perimeter security","  Incident investigation"]),
    ((0,50,100), "FORENSIC BOLO",
     ["Be On the Lookout search","Search hours of video in","minutes - not days","Identify persons of interest","","Use cases:","  Find suspect across feeds","  Locate missing persons","  Reconstruct incident timeline","  Legal evidence preparation"]),
    (MDBLUE, "MAVEN PARKING",
     ["Automated lot monitoring","Real-time occupancy tracking","Entry/exit event detection","Revenue analytics","","Use cases:","  Commercial parking ops","  Campus & facility mgmt","  Revenue optimization","  Overflow & capacity alerts"]),
    (LTBLUE, "PPE & SAFETY",
     ["Personal protective equipment","detection via Moondream VLM","Hard hats, vests, gloves,","goggles and more","","Use cases:","  Construction site safety","  Industrial compliance","  OSHA audit support","  Real-time safety alerts"]),
]
for i,(bgc,ttl,buls) in enumerate(modules):
    card(s, 0.2+i*3.27, 1.4, 3.1, 5.6, bgc, ttl, CYAN, buls, bullet_color=OFFWHT, title_size=13)

# SLIDE 11
s = slide()
content_frame(s, "Camera Health Monitoring", "Automated inspections - so you never miss a coverage gap")
rect(s, 0.2, 1.4, 7.0, 5.5, NAVY)
txb(s, [
    ("HOW IT WORKS", 14, True, CYAN, PP_ALIGN.LEFT),
    ("IntelliOptics inspects every camera view twice daily", 11, False, WHITE, PP_ALIGN.LEFT),
    ("Day and night comparison - confirms view consistency", 11, False, WHITE, PP_ALIGN.LEFT),
    ("Detects the slightest change in camera alignment", 11, False, WHITE, PP_ALIGN.LEFT),
    ("Maximum inspection rate: 10 frames per second", 11, False, WHITE, PP_ALIGN.LEFT),
    ("Historical health stats: trends, uptime, anomaly counts", 11, False, WHITE, PP_ALIGN.LEFT),
    ("Automated problem shortlists for cameras needing review", 11, False, WHITE, PP_ALIGN.LEFT),
    ("Custom report generation on camera health KPIs", 11, False, WHITE, PP_ALIGN.LEFT),
    ("Instant alerts when degradation or misalignment detected", 11, False, WHITE, PP_ALIGN.LEFT),
], 0.35, 1.55, 6.7, 5.2)
for (x,y,bgc,lbl,desc) in [
    (7.5, 1.4,  (16,130,90),   "NORMAL",     "Full coverage confirmed. View matches reference frame."),
    (10.5,1.4,  (180,100,0),   "BLURRY",     "Lens obscured or out of focus. Maintenance required."),
    (7.5, 4.0,  (160,30,30),   "TOO BRIGHT", "Overexposed. Lighting change or IR issue detected."),
    (10.5,4.0,  (30,30,30),    "TOO DARK",   "Underexposed. Power, IR, or obstruction detected."),
]:
    rect(s, x, y, 2.7, 2.4, bgc)
    txb(s, [(lbl,16,True,WHITE,PP_ALIGN.CENTER),(desc,10,False,WHITE,PP_ALIGN.CENTER)],
        x+0.1, y+0.2, 2.5, 2.0)

# SLIDE 12
s = slide()
content_frame(s, "IntelliOptics Benefits", "Measurable value delivered from day one")
card(s, 0.2, 1.4, 4.2, 5.5, DKBLUE, "COST SAVINGS", CYAN,
    ["Eliminates manual monitoring labor costs","Prevents unnecessary site visits and truck rolls",
     "Catch camera failures before incidents, not after","Proven: $57,000+ saved by St. Louis PD over 3 years",
     "Avoid litigation costs from undocumented gaps","Reduce insurance exposure with verified coverage"],
    bullet_color=OFFWHT)
card(s, 4.6, 1.4, 4.2, 5.5, BLUE, "SEAMLESS AI INTEGRATION", WHITE,
    ["AI analytics layer on your existing camera systems","No rip-and-replace - works with any IP camera brand",
     "Live in under one day - single install script","Unified cloud + edge Docker stack",
     "170+ API endpoints for custom integrations","Scales from single site to enterprise-wide"])
card(s, 9.0, 1.4, 4.2, 5.5, MDBLUE, "SECURITY ENHANCEMENT", CYAN,
    ["Real-time alerts on incidents and anomalies","Proactive vs reactive security posture",
     "SOC 2 compliant architecture","AES-256 data encryption at rest and in transit",
     "Full audit trail + role-based access logs","Forensic search for rapid incident response"],
    bullet_color=OFFWHT)

# SLIDE 13
s = slide()
content_frame(s, "IntelliOptics Time to Value", "Operational from day one - not months later")
rect(s, 0.2, 1.4, 6.3, 5.5, (255,240,240))
txb(s, [
    ("TRADITIONAL COMPUTER VISION", 14, True, (160,30,30), PP_ALIGN.LEFT),
    ("Day 1: Months of data collection required", 11, False, (100,30,30), PP_ALIGN.LEFT),
    ("Day 30-90: Model training begins", 11, False, (100,30,30), PP_ALIGN.LEFT),
    ("Day 90-180: Model training completes", 11, False, (100,30,30), PP_ALIGN.LEFT),
    ("Day 180+: Deploy - and hope it works", 11, False, (100,30,30), PP_ALIGN.LEFT),
    ("", 8, False, (100,30,30), PP_ALIGN.LEFT),
    ("Requires: ML engineers, labeled datasets,", 11, False, (100,30,30), PP_ALIGN.LEFT),
    ("GPU clusters, and dedicated model ops teams", 11, False, (100,30,30), PP_ALIGN.LEFT),
    ("", 8, False, (100,30,30), PP_ALIGN.LEFT),
    ("Result: Brittle models that fail on anything", 11, False, (160,30,30), PP_ALIGN.LEFT),
    ("outside their training distribution", 11, False, (160,30,30), PP_ALIGN.LEFT),
], 0.35, 1.55, 6.0, 5.2)
rect(s, 6.8, 1.4, 6.3, 5.5, (230,255,240))
txb(s, [
    ("INTELLIOPTICS 2.5", 14, True, (0,100,50), PP_ALIGN.LEFT),
    ("Day 1, Hour 1: Install, connect cameras, go live", 11, False, (0,60,30), PP_ALIGN.LEFT),
    ("Day 1, Hour 1: YOLOE detecting 376 classes", 11, False, (0,60,30), PP_ALIGN.LEFT),
    ("Day 1, Hour 1: VLM answering natural language queries", 11, False, (0,60,30), PP_ALIGN.LEFT),
    ("Day 1: Camera health baseline established", 11, False, (0,60,30), PP_ALIGN.LEFT),
    ("", 8, False, (0,60,30), PP_ALIGN.LEFT),
    ("Requires: One PowerShell install script.", 11, False, (0,60,30), PP_ALIGN.LEFT),
    ("That is it.", 11, True, (0,100,50), PP_ALIGN.LEFT),
    ("", 8, False, (0,60,30), PP_ALIGN.LEFT),
    ("Result: Immediate AI coverage across all cameras.", 11, False, (0,100,50), PP_ALIGN.LEFT),
    ("No training. No waiting. Zero ML expertise needed.", 11, True, (0,100,50), PP_ALIGN.LEFT),
], 6.95, 1.55, 6.0, 5.2)

# SLIDE 14
s = slide()
content_frame(s, "IntelliOptics - A Closer Look", "From raw footage to actionable intelligence in real time")
card(s, 0.2, 1.4, 4.2, 5.5, NAVY, "MONITORING & ALERTING", CYAN,
    ["Live 30fps bounding box overlay","Real-time feed health status dashboard",
     "Instant alerts: SMS, email, webhook","Per-camera and per-zone detection rules",
     "Confidence scores on all AI detections","VLM natural language answers inline",
     "Webcam + YouTube live capture support","Session-managed frame caching"],
    bullet_color=OFFWHT)
card(s, 4.6, 1.4, 4.2, 5.5, BLUE, "DRIVE AUTOMATION", WHITE,
    ["Gate open/close logic automation","PPE compliance auto-flagging and alerts",
     "Parking lot occupancy trigger events","BOLO match triggers immediate notification",
     "Vehicle access control rule enforcement","Custom AI rule builder - no code required",
     "Workflow integration via webhooks","Idempotent event processing pipeline"])
card(s, 9.0, 1.4, 4.2, 5.5, MDBLUE, "INSIGHTS & REPORTING", CYAN,
    ["Historical camera health trend analytics","Detection frequency and pattern reports",
     "Per-location performance metrics","Automated problem shortlist generation",
     "Scheduled custom report delivery","Compliance audit logs and exports",
     "Admin dashboard with full audit trail","Export to PDF and CSV formats"],
    bullet_color=OFFWHT)

# SLIDE 15
s = slide()
dark_frame(s)
rect(s, 0.8, 1.4, 11.7, 2.7, DKBLUE)
tx(s, '"IntelliOptics delivered measurable ROI before year one was over.\nThe automated camera health monitoring alone replaced two FTE review hours per day."',
   1.0, 1.6, 11.3, 2.3, size=17, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

tx(s, "$57,000", 0.8, 4.4, 5.8, 1.3, size=52, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
tx(s, "in reduced labor and equipment costs over 3 years", 0.8, 5.6, 5.8, 0.6, size=14, color=WHITE, align=PP_ALIGN.CENTER)
tx(s, "- ST. LOUIS POLICE DEPARTMENT", 0.8, 6.2, 5.8, 0.5, size=12, bold=True, color=(180,200,230), align=PP_ALIGN.CENTER)
for y,bgc,lbl in [(4.3,BLUE,"376 DETECTION CLASSES"),(5.35,MDBLUE,"1-DAY DEPLOYMENT"),(6.4,DKGRAY,"ZERO CAMERA REPLACEMENT NEEDED")]:
    rect(s, 7.3, y, 5.7, 0.8, bgc)
    tx(s, lbl, 7.35, y+0.1, 5.6, 0.6, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# SLIDE 16
s = slide()
content_frame(s, "Enterprise-Ready Platform", "Production-grade infrastructure - secure, scalable, deployable anywhere")
rect(s, 0.2, 1.4, 6.5, 5.5, NAVY)
txb(s, [
    ("DEPLOYMENT & INFRASTRUCTURE", 14, True, CYAN, PP_ALIGN.LEFT),
    ("Single Docker network: intellioptics-net", 11, False, WHITE, PP_ALIGN.LEFT),
    ("7 services: cloud-nginx, backend, frontend, worker,", 11, False, WHITE, PP_ALIGN.LEFT),
    ("  edge-nginx, edge-api, edge-inference", 11, False, WHITE, PP_ALIGN.LEFT),
    ("One-command install via PowerShell script", 11, False, WHITE, PP_ALIGN.LEFT),
    ("Auto-migration from IntelliOptics 2.0 on first run", 11, False, WHITE, PP_ALIGN.LEFT),
    ("GPU opt-in: uncomment deploy block for CC 7.0+ GPUs", 11, False, WHITE, PP_ALIGN.LEFT),
    ("VLM pre-loaded at startup - zero first-request cold start", 11, False, WHITE, PP_ALIGN.LEFT),
    ("Supabase PostgreSQL + Storage backend", 11, False, WHITE, PP_ALIGN.LEFT),
    ("SendGrid email + Twilio SMS alert delivery", 11, False, WHITE, PP_ALIGN.LEFT),
], 0.35, 1.55, 6.2, 5.2)
rect(s, 7.0, 1.4, 6.1, 5.5, OFFWHT)
txb(s, [
    ("SCALE & SECURITY", 14, True, BLUE, PP_ALIGN.LEFT),
    ("SOC 2 compliant architecture", 11, False, DKGRAY, PP_ALIGN.LEFT),
    ("AES-256-GCM encryption for all sensitive data", 11, False, DKGRAY, PP_ALIGN.LEFT),
    ("Role-based access: admin, reviewer, operator", 11, False, DKGRAY, PP_ALIGN.LEFT),
    ("Full audit event logging - tamper-evident records", 11, False, DKGRAY, PP_ALIGN.LEFT),
    ("17 frontend pages across full workflow", 11, False, DKGRAY, PP_ALIGN.LEFT),
    ("19 backend API routers", 11, False, DKGRAY, PP_ALIGN.LEFT),
    ("170+ REST endpoints - deny-by-default RBAC", 11, False, DKGRAY, PP_ALIGN.LEFT),
    ("Multi-site: cloud + edge on unified Docker stack", 11, False, DKGRAY, PP_ALIGN.LEFT),
    ("Scales horizontally - add edge nodes as needed", 11, False, DKGRAY, PP_ALIGN.LEFT),
], 7.15, 1.55, 5.8, 5.2)

# SLIDE 17
s = slide()
dark_frame(s)
tx(s, "Our Promise", 0.5, 0.7, 12, 1.1, size=40, bold=True, color=WHITE)
tx(s, "We make sure your surveillance is working - when it matters most.", 0.5, 1.9, 12, 0.7, size=18, italic=True, color=CYAN)
rect(s, 0.5, 2.7, 12.5, 0.05, BLUE)
promises_l = [
    "Real-time monitoring of every camera feed",
    "Alerted the moment a feed drops or misaligns",
    "Never lose coverage without knowing",
]
promises_r = [
    "Early warnings before critical moments are missed",
    "Support without burdening your IT team",
    "Peace of mind during audits, investigations & incidents",
]
for i,p in enumerate(promises_l):
    tx(s, f"  {p}", 0.5, 3.0+i*0.9, 6.2, 0.75, size=14, color=WHITE)
for i,p in enumerate(promises_r):
    tx(s, f"  {p}", 7.0, 3.0+i*0.9, 6.1, 0.75, size=14, color=WHITE)
rect(s, 0.5, 6.2, 12.5, 0.07, BLUE)
tx(s, "IntelliOptics 2.5  -  Your cameras, comprehended.", 0.5, 6.35, 12.5, 0.5, size=13, italic=True, color=CYAN, align=PP_ALIGN.CENTER)

# SLIDE 18
s = slide()
dark_frame(s)
tx(s, "Thank You", 0.5, 1.2, 12.3, 1.6, size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tx(s, "IntelliOptics 2.5  -  AI-Powered Visual Intelligence", 0.5, 2.9, 12.3, 0.8, size=20, color=CYAN, align=PP_ALIGN.CENTER)
rect(s, 3.2, 4.0, 6.9, 2.7, DKBLUE)
txb(s, [
    ("Haile Hantal", 18, True, WHITE, PP_ALIGN.CENTER),
    ("Co-Founder, CXO", 14, False, CYAN, PP_ALIGN.CENTER),
    ("314-760-7007", 13, False, WHITE, PP_ALIGN.CENTER),
    ("admin@4wardmotions.com", 13, False, WHITE, PP_ALIGN.CENTER),
    ("www.intellioptics.com", 13, False, CYAN, PP_ALIGN.CENTER),
], 3.3, 4.15, 6.7, 2.4)

prs.save(OUT)
print(f"SUCCESS: {len(prs.slides)} slides saved to {OUT}")
